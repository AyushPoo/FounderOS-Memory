#!/usr/bin/env python3
"""
Style Extractor — Founder Systems
Extracts primary colors and layout hints from uploaded PDF or PPTX reference decks.
Outputs a style override JSON compatible with generate_deck.py custom style.

Usage:
  python3 extract_style.py --input /tmp/reference.pdf --output /tmp/extracted_style.json
  python3 extract_style.py --input /tmp/reference.pptx --output /tmp/extracted_style.json
"""

import json
import sys
import argparse
from pathlib import Path
from collections import Counter


def extract_from_pdf(path: str) -> dict:
    """Extract dominant colors from first 2 pages of a PDF."""
    colors = {"bg": "#FFFFFF", "accent": "#1A1A1A", "text": "#111111"}
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        page = doc[0]

        # Get text to infer text color
        blocks = page.get_text("dict").get("blocks", [])
        text_colors = []
        for block in blocks:
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    c = span.get("color", 0)
                    if c:
                        r = (c >> 16) & 0xFF
                        g = (c >> 8) & 0xFF
                        b = c & 0xFF
                        text_colors.append((r, g, b))

        # Get background via pixmap sampling
        mat = fitz.Matrix(0.5, 0.5)
        pix = page.get_pixmap(matrix=mat)
        img_data = pix.samples
        w, h = pix.width, pix.height

        # Sample corners for background color
        corner_colors = []
        for sample_x in [0, w - 1]:
            for sample_y in [0, h - 1]:
                idx = (sample_y * w + sample_x) * pix.n
                if pix.n >= 3:
                    corner_colors.append((img_data[idx], img_data[idx + 1], img_data[idx + 2]))

        if corner_colors:
            avg_r = sum(c[0] for c in corner_colors) // len(corner_colors)
            avg_g = sum(c[1] for c in corner_colors) // len(corner_colors)
            avg_b = sum(c[2] for c in corner_colors) // len(corner_colors)
            colors["bg"] = f"#{avg_r:02X}{avg_g:02X}{avg_b:02X}"

        if text_colors:
            most_common_text = Counter(text_colors).most_common(1)[0][0]
            r, g, b = most_common_text
            colors["text"] = f"#{r:02X}{g:02X}{b:02X}"

        # Scan for accent: look for highly saturated colors in the center of the page
        center_sample = []
        cx, cy = w // 2, h // 2
        for x in range(max(0, cx - 50), min(w, cx + 50)):
            for y in range(max(0, cy - 50), min(h, cy + 50)):
                idx = (y * w + x) * pix.n
                if pix.n >= 3:
                    r2, g2, b2 = img_data[idx], img_data[idx + 1], img_data[idx + 2]
                    # Check saturation
                    max_c = max(r2, g2, b2)
                    min_c = min(r2, g2, b2)
                    if max_c > 50 and max_c - min_c > 60:  # Saturated
                        center_sample.append((r2, g2, b2))

        if center_sample:
            most_common_accent = Counter(center_sample).most_common(1)[0][0]
            r, g, b = most_common_accent
            colors["accent"] = f"#{r:02X}{g:02X}{b:02X}"

        doc.close()

    except ImportError:
        print("PyMuPDF not installed, using default colors. Run: pip3 install PyMuPDF", file=sys.stderr)
    except Exception as e:
        print(f"PDF extraction error: {e}", file=sys.stderr)

    return colors


def extract_from_pptx(path: str) -> dict:
    """Extract theme colors from a PPTX file."""
    colors = {"bg": "#FFFFFF", "accent": "#1A1A1A", "text": "#111111"}
    try:
        from pptx import Presentation
        from pptx.util import Pt
        import lxml.etree as etree

        prs = Presentation(path)
        theme_colors = []

        # Extract from slide master theme
        for layout in prs.slide_layouts:
            theme = layout.slide_master.theme_color_map
            break

        # Get slide 1 background + shapes for color hints
        if prs.slides:
            slide = prs.slides[0]

            # Background fill
            bg_fill = slide.background.fill
            if bg_fill.type is not None:
                try:
                    bg_rgb = bg_fill.fore_color.rgb
                    colors["bg"] = f"#{bg_rgb}"
                except Exception:
                    pass

            # Text and accent from shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                rgb = run.font.color.rgb
                                theme_colors.append(str(rgb))
                            except Exception:
                                pass

        if theme_colors:
            most_common = Counter(theme_colors).most_common(1)[0][0]
            r = int(most_common[0:2], 16)
            g = int(most_common[2:4], 16)
            b = int(most_common[4:6], 16)
            # If it's a bright saturated color, use as accent; else as text
            max_c = max(r, g, b)
            min_c = min(r, g, b)
            if max_c - min_c > 50:
                colors["accent"] = f"#{most_common}"
            else:
                colors["text"] = f"#{most_common}"

    except ImportError:
        print("python-pptx not installed. Run: pip3 install python-pptx", file=sys.stderr)
    except Exception as e:
        print(f"PPTX extraction error: {e}", file=sys.stderr)

    return colors


def is_dark_color(hex_color: str) -> bool:
    """Return True if a hex color is dark (luminance < 0.4)."""
    h = hex_color.lstrip("#")
    if len(h) < 6:
        return False
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    return luminance < 0.4


def build_custom_style(extracted: dict) -> dict:
    """Build a full style config from extracted colors."""
    bg = extracted.get("bg", "#FFFFFF")
    accent = extracted.get("accent", "#007AFF")
    dark_mode = is_dark_color(bg)

    if dark_mode:
        return {
            "name": "Custom (Dark)",
            "bg_slide": bg,
            "bg_cover": bg,
            "accent": accent,
            "accent2": accent,
            "text_on_accent": "#FFFFFF",
            "text_primary": "#F0F0F0",
            "text_secondary": "#B0B0B0",
            "text_muted": "#707070",
            "divider": "#222222",
            "card_bg": "#1A1A1A",
            "card_border": "#2A2A2A",
            "font": "'Inter', 'Helvetica Neue', Arial, sans-serif",
            "cover_style": "dark_gradient",
        }
    else:
        return {
            "name": "Custom (Light)",
            "bg_slide": bg,
            "bg_cover": accent,
            "accent": accent,
            "text_on_accent": "#FFFFFF",
            "text_primary": "#111111",
            "text_secondary": "#555555",
            "text_muted": "#888888",
            "divider": "#E5E5E5",
            "card_bg": "#F5F5F5",
            "card_border": "#E0E0E0",
            "font": "'Inter', 'Helvetica Neue', Arial, sans-serif",
            "cover_style": "orange_bold",
        }


def extract_style(input_path: str, output_path: str) -> bool:
    """Main entry point."""
    p = Path(input_path)
    if not p.exists():
        print(f"ERROR: File not found: {input_path}", file=sys.stderr)
        return False

    ext = p.suffix.lower()
    print(f"Extracting style from: {p.name}")

    if ext == ".pdf":
        extracted = extract_from_pdf(input_path)
    elif ext in (".pptx", ".ppt"):
        extracted = extract_from_pptx(input_path)
    else:
        print(f"ERROR: Unsupported file type: {ext}. Use .pdf or .pptx", file=sys.stderr)
        return False

    custom_style = build_custom_style(extracted)
    print(f"Detected: bg={extracted['bg']}, accent={extracted['accent']}, dark={is_dark_color(extracted['bg'])}")

    with open(output_path, "w") as f:
        json.dump(custom_style, f, indent=2)
    print(f"Style saved: {output_path}")
    return True


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract style from reference deck")
    parser.add_argument("--input", required=True, help="Path to reference PDF or PPTX")
    parser.add_argument("--output", required=True, help="Path for output style JSON")
    args = parser.parse_args()

    success = extract_style(args.input, args.output)
    sys.exit(0 if success else 1)
